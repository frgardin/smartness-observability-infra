from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Paleta ───────────────────────────────────────────────────────────────────
DARK_BG    = RGBColor(0x1A, 0x1A, 0x2E)   # fundo escuro
ACCENT     = RGBColor(0xE9, 0x4F, 0x37)   # laranja/vermelho vibrante
ACCENT2    = RGBColor(0x39, 0x3E, 0x46)   # cinza escuro para painéis
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
YELLOW     = RGBColor(0xF5, 0xA6, 0x23)

AUTHOR  = "Felipe Rezende Gardin"
STUDENT_ID = "215766"

# ── Helpers ───────────────────────────────────────────────────────────────────
def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def title_slide(prs, title, subtitle, author, sid):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, DARK_BG)
    # barra lateral colorida
    add_rect(slide, 0, 0, 0.08, 7.5, ACCENT)
    # título
    add_text(slide, title, 0.4, 1.8, 9.0, 1.4, size=36, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    # linha decorativa
    add_rect(slide, 2.5, 3.35, 5.0, 0.04, ACCENT)
    # subtítulo
    add_text(slide, subtitle, 0.4, 3.5, 9.0, 0.8, size=20,
             color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    # autor e matrícula
    add_text(slide, f"{author}  •  Matrícula: {sid}", 0.4, 6.6, 9.0, 0.5,
             size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

def content_slide(prs, title, bullets):
    """bullets: lista de (texto, nivel) onde nivel 0=principal, 1=sub"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    # barra topo
    add_rect(slide, 0, 0, 10, 1.1, ACCENT2)
    add_rect(slide, 0, 0, 0.08, 7.5, ACCENT)
    # título
    add_text(slide, title, 0.25, 0.12, 9.4, 0.85, size=24, bold=True,
             color=YELLOW, align=PP_ALIGN.LEFT)

    txb = slide.shapes.add_textbox(Inches(0.4), Inches(1.3), Inches(9.2), Inches(5.8))
    tf = txb.text_frame
    tf.word_wrap = True

    first = True
    for text, level in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.level = level
        p.space_before = Pt(4 if level == 0 else 2)
        run = p.add_run()
        bullet_char = "▸ " if level == 0 else "  – "
        run.text = bullet_char + text
        run.font.size = Pt(15 if level == 0 else 13)
        run.font.bold = (level == 0)
        run.font.color.rgb = WHITE if level == 0 else LIGHT_GRAY
    return slide

def code_slide(prs, title, code_text, note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BG)
    add_rect(slide, 0, 0, 10, 1.1, ACCENT2)
    add_rect(slide, 0, 0, 0.08, 7.5, ACCENT)
    add_text(slide, title, 0.25, 0.12, 9.4, 0.85, size=24, bold=True,
             color=YELLOW, align=PP_ALIGN.LEFT)
    # caixa de código
    add_rect(slide, 0.3, 1.25, 9.4, 5.2 if not note else 4.4, ACCENT2)
    txb = slide.shapes.add_textbox(Inches(0.5), Inches(1.35),
                                    Inches(9.0), Inches(5.0 if not note else 4.2))
    tf = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = code_text
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0xA8, 0xFF, 0x3E)  # verde terminal
    run.font.name = "Courier New"
    if note:
        add_text(slide, f"💡 {note}", 0.3, 5.85, 9.4, 0.6, size=13,
                 color=YELLOW)

# ═══════════════════════════════════════════════════════════════════════════════
#  APRESENTAÇÃO 1 — PROJETO
# ═══════════════════════════════════════════════════════════════════════════════
def build_projeto():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    title_slide(prs,
        "Smartness Observability\nInfrastructure",
        "Sistema de monitoramento distribuído com Prometheus, Grafana e Ansible",
        AUTHOR, STUDENT_ID)

    content_slide(prs, "O que é esse projeto", [
        ("Sistema de observabilidade centralizado para 40+ máquinas", 0),
        ("Uma máquina master coleta e visualiza métricas de todas as outras (slaves)", 1),
        ("Métricas em tempo real: CPU, memória, disco e rede", 1),
        ("Por que monitorar?", 0),
        ("Detectar problemas antes que virem incidentes", 1),
        ("Entender o comportamento real das máquinas sob carga", 1),
        ("Base para decisões de capacidade e escalabilidade", 1),
    ])

    content_slide(prs, "Arquitetura — Visão Geral", [
        ("Topologia Master / Slave", 0),
        ("Master: roda Prometheus + Grafana via Docker Compose", 1),
        ("Slaves: cada máquina expõe métricas via Node Exporter na porta 9100", 1),
        ("Modelo Pull — o Prometheus busca as métricas", 0),
        ("Slaves não precisam saber onde está o master", 1),
        ("Falha de conexão = alerta automático (target DOWN)", 1),
        ("Descoberta dinâmica via file_sd", 0),
        ("Adicionar uma máquina = editar um JSON + rodar Ansible", 1),
        ("Prometheus recarrega em até 30s sem restart", 1),
    ])

    code_slide(prs, "Arquitetura — Fluxo de Dados",
"""Master Node
  Prometheus (:9090)  ←── pull /metrics a cada 15s ───►  Slave 01: node_exporter :9100
  Grafana    (:3000)  ←── PromQL queries             ───►  Slave 02: node_exporter :9100
  file_sd/targets.json                                ───►  ...        40+ máquinas

Ansible (local)  ──── SSH ────►  instala node_exporter em todas as slaves
                              └►  atualiza file_sd/targets.json no master""",
        "O Prometheus é o único agente ativo — as slaves só respondem")

    content_slide(prs, "Stack de Tecnologias", [
        ("Prometheus — banco de séries temporais", 0),
        ("Coleta e armazena métricas de todas as máquinas", 1),
        ("Linguagem de consulta própria: PromQL", 1),
        ("Grafana — visualização", 0),
        ("Dashboards interativos com gráficos e tabelas", 1),
        ("Configurado inteiramente como código (provisionamento)", 1),
        ("Node Exporter — agente nas slaves", 0),
        ("Expõe métricas de SO via HTTP em texto simples", 1),
        ("Instalado como serviço systemd (sem Docker nas slaves)", 1),
        ("Ansible — automação de deploy", 0),
        ("Instala node-exporter em 40+ máquinas via SSH em paralelo", 1),
        ("Idempotente: pode rodar N vezes sem efeito colateral", 1),
    ])

    content_slide(prs, "Descoberta Dinâmica de Targets (file_sd)", [
        ("Problema: como o Prometheus sabe quais máquinas monitorar?", 0),
        ("Solução: file-based Service Discovery", 0),
        ("Prometheus lê um arquivo JSON com a lista de targets", 1),
        ("Recarrega automaticamente a cada 30 segundos — sem restart", 1),
        ("Label hostname → sobrescreve instance no Prometheus", 1),
        ("Grafana mostra 'maquina-01' em vez de '10.0.0.10:9100'", 1),
        ("Fluxo de adição de uma nova máquina", 0),
        ("1. Adicionar host no ansible/inventory/hosts.yml", 1),
        ("2. Rodar: ansible-playbook playbooks/deploy.yml", 1),
        ("3. Em até 30s a máquina aparece no Grafana automaticamente", 1),
    ])

    code_slide(prs, "Automação Ansible — O que acontece no deploy",
"""# Para cada máquina no inventário, o Ansible:

1. Cria usuário de sistema 'node_exporter' (sem shell, sem home)

2. Baixa o binário do GitHub:
   node_exporter-1.8.2.linux-amd64.tar.gz

3. Instala em /usr/local/bin/node_exporter

4. Cria /etc/systemd/system/node_exporter.service

5. systemctl enable --now node_exporter

6. Abre porta 9100 no firewall (ufw/firewalld)

7. Atualiza file_sd/targets.json no master com todos os hosts""",
        "Executa em paralelo em todas as máquinas — 40 hosts ≈ mesmo tempo que 1")

    content_slide(prs, "Dashboards — Fleet Overview", [
        ("Visão de toda a frota em uma tabela única", 0),
        ("Uma linha por máquina — sem precisar selecionar", 1),
        ("Ideal para detectar rapidamente qual máquina está com problema", 1),
        ("Colunas da tabela", 0),
        ("Máquina — nome do host", 1),
        ("Status — UP (verde) / DOWN (vermelho)", 1),
        ("CPU % — com coloração por threshold (verde/amarelo/vermelho)", 1),
        ("RAM Usada e RAM Total em bytes formatados automaticamente", 1),
        ("Disk % da partição raiz — com threshold visual", 1),
    ])

    content_slide(prs, "Dashboards — Node Metrics (por máquina)", [
        ("Overview — valores instantâneos no topo", 0),
        ("CPU %, RAM Usada, RAM Total, Disk Usado, Disk Total, Uptime", 1),
        ("CPU — gráficos temporais", 0),
        ("Uso % por modo: user, system, iowait (identifica gargalos)", 1),
        ("Load Average 1m/5m/15m: compara com número de cores", 1),
        ("Memory — evolução ao longo do tempo", 0),
        ("RAM: total, used, available, cached — em bytes absolutos", 1),
        ("Swap: indica falta real de memória quando sobe", 1),
        ("Disk & Network", 0),
        ("Espaço por partição e I/O (leitura/escrita em bytes/s)", 1),
        ("Throughput de rede por interface (RX/TX)", 1),
    ])

    content_slide(prs, "Conceitos Fundamentais Aprendidos", [
        ("Time Series multidimensional", 0),
        ("Cada métrica identificada por nome + labels (chave=valor)", 1),
        ("Labels permitem filtrar e agregar de qualquer ângulo depois", 1),
        ("PromQL — linguagem de consulta", 0),
        ("rate() para contadores, avg/sum/max para agregação", 1),
        ("Janela de tempo [1m] vs [5m]: reatividade vs suavização", 1),
        ("Pull model", 0),
        ("Prometheus é o agente ativo — slaves só respondem", 1),
        ("Falha de rede vira métrica (target DOWN)", 1),
        ("Idempotência no Ansible", 0),
        ("Rodar o playbook dez vezes = mesmo resultado que uma vez", 1),
    ])

    code_slide(prs, "Como rodar o projeto",
"""# 1. Subir o master
docker compose up -d

# 2. Acessar Grafana
http://localhost:3000   (admin / admin)

# 3. Adicionar slaves no inventário
vim ansible/inventory/hosts.yml

# 4. Deploy nas slaves
ansible-playbook ansible/playbooks/deploy.yml \\
  -i ansible/inventory/hosts.yml

# 5. Verificar quantas máquinas estão sendo monitoradas
curl -s http://localhost:9090/api/v1/targets | jq \\
  '.data.activeTargets | length'""")

    prs.save("docs/projeto.pptx")
    print("✓ docs/projeto.pptx gerado")


# ═══════════════════════════════════════════════════════════════════════════════
#  APRESENTAÇÃO 2 — IA NO DESENVOLVIMENTO
# ═══════════════════════════════════════════════════════════════════════════════
def build_ia():
    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    title_slide(prs,
        "Como a IA Auxiliou no\nDesenvolvimento",
        "O papel do Claude (Anthropic) na construção da infraestrutura de observabilidade",
        AUTHOR, STUDENT_ID)

    content_slide(prs, "Modelo de Colaboração", [
        ("Ensino iterativo — não automação cega", 0),
        ("Cada bloco de código foi acompanhado de explicação do porquê", 1),
        ("Nenhuma etapa foi avançada sem o desenvolvedor entender a anterior", 1),
        ("O ciclo que emergiu naturalmente", 0),
        ("1. Desenvolvedor declara intenção ou dúvida", 1),
        ("2. IA explica o conceito subjacente", 1),
        ("3. IA gera código/configuração com justificativa", 1),
        ("4. Desenvolvedor testa na prática", 1),
        ("5. Problema real aparece → IA diagnostica com causa raiz", 1),
        ("6. Correção aplicada + explicação do porquê", 1),
    ])

    content_slide(prs, "Etapa 1 — Explicação dos Fundamentos", [
        ("Antes de qualquer código: conceitos da documentação oficial", 0),
        ("Modelo de dados multidimensional", 1),
        ("O que são time series, labels e por que são mais poderosos que nomes estáticos", 1),
        ("PromQL — da query simples à agregação complexa", 1),
        ("Pull model vs Push model — quando cada um faz sentido", 1),
        ("Service discovery — static_configs vs file_sd", 1),
        ("Resultado", 0),
        ("Vocabulário e modelo mental estabelecidos antes da implementação", 1),
        ("O desenvolvedor sabia o que estava construindo, não só como rodar", 1),
    ])

    content_slide(prs, "Etapa 2 — Diagnóstico de Rede (Linux vs Mac)", [
        ("Problema real encontrado", 0),
        ("Node Exporter com network_mode: host não era acessível pelo Prometheus", 1),
        ("curl localhost:9100 funcionava no host, mas falhava de dentro do container", 1),
        ("Processo de diagnóstico da IA", 0),
        ("Testou conectividade no host → OK", 1),
        ("Testou de dentro do container Prometheus → timeout", 1),
        ("Identificou causa: Docker bridge no Linux não roteia para o host por padrão", 1),
        ("Diferente do Mac/Windows onde host.docker.internal funciona nativamente", 1),
        ("Solução aplicada", 0),
        ("Remover network_mode: host, expor porta normalmente, usar nome de serviço Docker", 1),
        ("node_exporter:9100 resolvido pelo DNS interno do Docker", 1),
    ])

    content_slide(prs, "Etapa 3 — Grafana como Código", [
        ("Provisionamento via arquivos (não configuração manual pela UI)", 0),
        ("Repeatability: qualquer clone do repositório tem o mesmo Grafana", 1),
        ("Versionamento: mudanças em dashboards ficam no git", 1),
        ("Problema de UID do datasource", 0),
        ("Primeira tentativa usou uid: '-- Default --' no JSON do dashboard", 1),
        ("Erro: Datasource provisioning error: data source not found", 1),
        ("Causa: Grafana tinha datasource antigo no volume sem UID definido", 1),
        ("Solução: UID explícito no datasource + docker compose down -v", 1),
        ("O que foi aprendido", 0),
        ("Como o Grafana reconcilia estado provisionado com banco SQLite interno", 1),
        ("Por que volumes Docker persistentes podem causar conflitos de estado", 1),
    ])

    content_slide(prs, "Etapa 4 — Diagnóstico de Query Incorreta de CPU", [
        ("Teste prático com stress-ng para verificar 100% de CPU", 0),
        ("Resultado esperado: gráfico mostrando 100%", 1),
        ("Resultado obtido: gráfico mostrando ~20%", 1),
        ("Load Average reagia instantaneamente, CPU % não", 1),
        ("Diagnóstico da IA", 0),
        ("rate(node_cpu_seconds_total[5m]) calcula média dos últimos 5 minutos", 1),
        ("Stress rodando 1min dentro de janela de 5min → 1/5 = 20%", 1),
        ("Load Average é calculado pelo kernel — reage em segundos", 1),
        ("Solução: trocar [5m] por [1m]", 0),
        ("Mais reativo: após 60s de stress, gráfico reflete 100%", 1),
        ("Tradeoff explicado: [1m] mostra mais ruído, [5m] suaviza demais", 1),
        ("[1m] é o padrão da indústria para CPU", 1),
    ])

    content_slide(prs, "Etapa 5 — Automação Ansible", [
        ("Necessidade: instalar node-exporter em 40+ máquinas", 0),
        ("A IA recomendou Ansible com justificativa", 0),
        ("Agentless — só precisa de SSH, nada a instalar nas slaves", 1),
        ("Idempotente — pode reexecutar sem efeito colateral", 1),
        ("Integração natural com file_sd: inventário vira fonte de verdade", 1),
        ("Design da solução", 0),
        ("Role node_exporter: usuário de sistema, binário, systemd, firewall", 1),
        ("Template Jinja2 transforma inventário em targets.json automaticamente", 1),
        ("relabel_configs para mostrar hostnames no Grafana em vez de IPs", 1),
        ("Resultado", 0),
        ("Deploy de 40 máquinas em paralelo com um único comando", 1),
        ("Grafana atualiza automaticamente sem nenhuma configuração manual", 1),
    ])

    content_slide(prs, "O que foi decidido pelo desenvolvedor", [
        ("Escopo do projeto", 0),
        ("Monitorar 40+ máquinas com as métricas mais relevantes para operação", 1),
        ("Tecnologia principal: Prometheus + Grafana", 0),
        ("Escolha feita antes da conversa com a IA", 1),
        ("Abordagem de aprendizado", 0),
        ("Entender cada peça antes de avançar para a próxima", 1),
        ("Instalação nas slaves: binário + systemd (sem Docker)", 0),
        ("Escolhido por ser mais leve e não exigir Docker em cada slave", 1),
        ("Aceitação ou rejeição de sugestões da IA", 0),
        ("Toda sugestão foi avaliada — a IA apresentou opções com tradeoffs", 1),
    ])

    content_slide(prs, "O que a IA não fez", [
        ("Não tomou decisões de arquitetura sozinha", 0),
        ("Cada escolha foi apresentada com opções e tradeoffs explicados", 1),
        ("O desenvolvedor sempre escolheu a direção", 1),
        ("Não pulou etapas de aprendizado", 0),
        ("Quando havia dúvida conceitual, a IA ensinou antes de gerar código", 1),
        ("Não ignorou erros reais", 0),
        ("Todos os problemas foram diagnosticados com causa raiz", 1),
        ("A correção nunca veio sem a explicação do porquê", 1),
        ("Não gerou código que o desenvolvedor não entendesse", 0),
        ("Cada arquivo criado foi explicado: o quê, por quê, como funciona", 1),
    ])

    content_slide(prs, "Reflexão: IA como Ferramenta de Ensino", [
        ("Uso tradicional de IA em projetos técnicos", 0),
        ("'Me gere o código X' → copiar → colar → funciona (mas não aprende)", 1),
        ("O que foi feito aqui", 0),
        ("'Me explique X, depois me ajude a construir X' → aprende + constrói", 1),
        ("A IA atuou como um colega sênior disponível 24h", 0),
        ("Explica sem julgamento, diagnostica pacientemente, adapta ao ritmo", 1),
        ("Não substitui o raciocínio — amplifica a capacidade de aprender", 1),
        ("O resultado", 0),
        ("Infraestrutura funcionando E entendimento de por que funciona", 1),
        ("Capacidade de manter, evoluir e ensinar o projeto a outros", 1),
    ])

    prs.save("docs/ia-no-desenvolvimento.pptx")
    print("✓ docs/ia-no-desenvolvimento.pptx gerado")


if __name__ == "__main__":
    build_projeto()
    build_ia()
